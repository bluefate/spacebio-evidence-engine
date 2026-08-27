#!/usr/bin/env python3
"""Build final showcase PowerPoint artifacts for issue #195."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FINAL = ROOT / "docs" / "final"
SHOTS = FINAL / "screenshots"
LOGO = ROOT / "docs" / "brand" / "logo-wordmark.png"
HERO = ROOT / "docs" / "brand" / "hero-atmosphere.png"

NAVY = RGBColor(4, 24, 57)
BLUE = RGBColor(0, 111, 145)
CYAN = RGBColor(39, 185, 205)
LIGHT = RGBColor(239, 248, 252)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(73, 91, 116)
GOLD = RGBColor(220, 160, 40)


def add_bg(slide, *, dark: bool = False) -> None:
    color = NAVY if dark else WHITE
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color
    if dark and HERO.is_file():
        slide.shapes.add_picture(str(HERO), Inches(6.6), 0, width=Inches(6.733), height=Inches(7.5))
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.3), 0, Inches(7.1), Inches(7.5))
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = NAVY
        overlay.fill.transparency = 32
        overlay.line.fill.background()


def add_logo(slide, *, x: float = 0.55, y: float = 0.35, width: float = 2.55) -> None:
    slide.shapes.add_picture(str(LOGO), Inches(x), Inches(y), width=Inches(width))


def add_title(slide, title: str, subtitle: str | None = None, *, dark: bool = False) -> None:
    color = WHITE if dark else NAVY
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.9), Inches(12.0), Inches(0.75))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Aptos Display"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = color
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.72), Inches(1.65), Inches(11.8), Inches(0.6))
        p = sub.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.name = "Aptos"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(210, 226, 238) if dark else GRAY


def add_footer(slide, number: int, *, dark: bool = False) -> None:
    color = RGBColor(180, 210, 225) if dark else GRAY
    box = slide.shapes.add_textbox(Inches(0.7), Inches(7.1), Inches(12), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = f"Space Biology Evidence Engine  |  John Hernandez  |  {number}"
    p.font.name = "Aptos"
    p.font.size = Pt(8)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.RIGHT


def add_bullets(slide, bullets: list[str], *, x: float, y: float, w: float, h: float, dark: bool = False, size: int = 19) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    color = WHITE if dark else NAVY
    for i, text in enumerate(bullets):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = text
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(12)
        p.text = f"•  {text}"


def add_card(slide, title: str, body: str, *, x: float, y: float, w: float, h: float, accent: RGBColor = CYAN) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.14)
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Aptos Display"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p = tf.add_paragraph()
    p.text = body
    p.font.name = "Aptos"
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    p.space_before = Pt(5)


def add_screenshot(slide, filename: str, *, x: float, y: float, w: float, h: float) -> None:
    border = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x - 0.05), Inches(y - 0.05), Inches(w + 0.1), Inches(h + 0.1))
    border.fill.solid()
    border.fill.fore_color.rgb = WHITE
    border.line.color.rgb = CYAN
    border.line.width = Pt(1.5)
    slide.shapes.add_picture(str(SHOTS / filename), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def add_notes(slide, notes: str) -> None:
    frame = slide.notes_slide.notes_text_frame
    frame.text = notes


def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def build_pitch() -> None:
    prs = new_deck()

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, dark=True)
    add_logo(slide, x=0.75, y=0.55, width=3.6)
    title = slide.shapes.add_textbox(Inches(0.85), Inches(2.0), Inches(7.0), Inches(1.5))
    p = title.text_frame.paragraphs[0]
    p.text = "Space Biology\nEvidence Engine"
    p.font.name = "Aptos Display"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = WHITE
    sub = slide.shapes.add_textbox(Inches(0.9), Inches(3.72), Inches(6.0), Inches(1.0))
    p = sub.text_frame.paragraphs[0]
    p.text = "Citation-first research synthesis from a controlled corpus"
    p.font.name = "Aptos"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(190, 226, 235)
    by = slide.shapes.add_textbox(Inches(0.9), Inches(5.8), Inches(5.8), Inches(0.8))
    p = by.text_frame.paragraphs[0]
    p.text = "John Hernandez  |  jherna65@fau.edu\nFAU AI HootCamp  |  Final Project 2026"
    p.font.name = "Aptos"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    add_notes(slide, "Open with the challenge: researchers need trustworthy answers, not another general-purpose chatbot. Introduce the project as a local-first, citation-first evidence workspace.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "The problem", "Space biology evidence is fragmented and difficult to synthesize safely")
    add_card(slide, "Scattered evidence", "Relevant studies span publishers, organisms, experimental models, and terminology.", x=0.8, y=2.4, w=3.7, h=2.1)
    add_card(slide, "High trust bar", "Scientific users need passage-level provenance, not unsupported model recall.", x=4.8, y=2.4, w=3.7, h=2.1, accent=GOLD)
    add_card(slide, "Slow comparison", "Researchers repeatedly search, read, and reconcile study context by hand.", x=8.8, y=2.4, w=3.7, h=2.1, accent=BLUE)
    add_bullets(slide, ["Primary users: space-biology researchers, students, educators, and corpus maintainers"], x=1.1, y=5.05, w=11.0, h=0.8, size=18)
    add_footer(slide, 2)
    add_notes(slide, "Explain that the core problem is not access to a language model. It is evidence discovery with scientific provenance and explicit limits.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, dark=True)
    add_title(slide, "The solution", "A controlled, provenance-preserving evidence engine", dark=True)
    add_bullets(slide, ["23 owner-approved open-access publications", "Semantic + keyword retrieval over indexed passages", "Grounded answers with passage-level citations", "Explicit insufficient-evidence responses", "Local-first operation with optional OpenAI or Ollama generation"], x=0.9, y=2.3, w=5.3, h=3.8, dark=True, size=20)
    add_screenshot(slide, "01-home.png", x=6.55, y=2.05, w=6.1, h=4.24)
    add_footer(slide, 3, dark=True)
    add_notes(slide, "Walk through the controlled corpus and fail-closed approach. Stress that the system does not silently fall back to general model knowledge.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Architecture", "Retrieval and generation remain independently testable")
    stages = [("Approved PDFs", "License + quality checks"), ("Ingest", "Extract → chunk → embed"), ("Postgres + pgvector", "Metadata + vectors"), ("Hybrid retrieval", "Semantic + full-text"), ("Grounded answer", "Citation validation")]
    for i, (name, body) in enumerate(stages):
        x = 0.55 + i * 2.55
        add_card(slide, name, body, x=x, y=2.55, w=2.15, h=1.65, accent=CYAN if i < 4 else GOLD)
        if i < len(stages) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.17), Inches(3.12), Inches(0.32), Inches(0.35))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = BLUE; arrow.line.fill.background()
    add_bullets(slide, ["Next.js frontend", "FastAPI API", "SQLAlchemy + Alembic", "Local embeddings", "Provider-neutral LLM interface"], x=0.95, y=4.95, w=11.5, h=1.2, size=16)
    add_footer(slide, 4)
    add_notes(slide, "Describe the data flow left to right. The same provenance fields survive extraction, retrieval, answer generation, and citation display.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Grounded question answering", "Answers are constrained to retrieved corpus passages")
    add_screenshot(slide, "02-ask.png", x=0.65, y=2.05, w=7.0, h=4.86)
    add_card(slide, "Evidence sufficiency", "Weak or empty retrieval returns a fixed insufficient-evidence response without invoking the LLM.", x=8.05, y=2.25, w=4.5, h=1.55)
    add_card(slide, "Citation validation", "Unknown, missing, or invented citation markers fail closed instead of becoming successful answers.", x=8.05, y=4.05, w=4.5, h=1.55, accent=GOLD)
    add_footer(slide, 5)
    add_notes(slide, "Demonstrate the Ask form. Explain that every answer claim should map to a retrieved citation and that unsupported generation is rejected.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, dark=True)
    add_title(slide, "Controlled corpus and provenance", "Inventory visibility makes the evidence boundary explicit", dark=True)
    add_screenshot(slide, "04-corpus.png", x=0.65, y=1.95, w=7.2, h=5.0)
    add_screenshot(slide, "06-publication.png", x=8.15, y=1.95, w=4.55, h=5.0)
    add_footer(slide, 6, dark=True)
    add_notes(slide, "Show that users can inspect the approved corpus and open individual publication metadata, DOI, license, organism, exposure, and ingest status.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Search and study comparison", "Discovery without collapsing distinct evidence types")
    add_screenshot(slide, "03-search.png", x=0.55, y=2.0, w=6.05, h=4.2)
    add_screenshot(slide, "05-compare.png", x=6.75, y=2.0, w=6.05, h=4.2)
    add_bullets(slide, ["Search preserves IDs, source links, sections, and pages", "Compare keeps organism and experimental-system labels visible"], x=0.9, y=6.35, w=11.6, h=0.6, size=13)
    add_footer(slide, 7)
    add_notes(slide, "Explain semantic/keyword search and the compare workspace. Comparison is metadata-based and deliberately does not invent cross-study findings.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "AI capabilities and safeguards", "Useful automation with explicit scientific boundaries")
    capabilities = [("Local embeddings", "Sentence Transformers for private, low-cost semantic retrieval."), ("Provider-neutral generation", "OpenAI and local Ollama paths behind a shared interface."), ("Hybrid retrieval", "Vector similarity and PostgreSQL full-text retrieval with optional reranking."), ("Evidence gates", "Minimum passage and publication support before generation."), ("Citation-first output", "Passage IDs and publication provenance remain inspectable."), ("Hallucination checks", "Offline fixtures flag unsupported claims and citation mismatches.")]
    for i, (name, body) in enumerate(capabilities):
        add_card(slide, name, body, x=0.65 + (i % 3) * 4.22, y=2.05 + (i // 3) * 2.1, w=3.85, h=1.75, accent=[CYAN, BLUE, GOLD][i % 3])
    add_footer(slide, 8)
    add_notes(slide, "Cover the AI stack and safeguards. Keep the distinction clear: retrieval identifies evidence; generation summarizes only what retrieval supports.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, dark=True)
    add_title(slide, "Results and engineering quality", "A tested MVP with transparent limitations", dark=True)
    add_card(slide, "23 publications", "Approved controlled corpus inventory", x=0.75, y=2.2, w=3.55, h=1.55)
    add_card(slide, "292+ Python tests", "Unit, API, retrieval, ingestion, citation, and accessibility coverage", x=4.55, y=2.2, w=3.55, h=1.55, accent=GOLD)
    add_card(slide, "CI quality gates", "Ruff, Pyright, Python tests, Node typecheck and UI tests", x=8.35, y=2.2, w=3.55, h=1.55, accent=BLUE)
    add_bullets(slide, ["Citation correctness, hallucination, and retrieval evaluation harnesses", "No secrets committed; `.env` loading and fail-closed provider configuration", "Accessible interface patterns and keyboard-tested UI components"], x=1.0, y=4.35, w=11.2, h=1.8, dark=True, size=18)
    add_footer(slide, 9, dark=True)
    add_notes(slide, "Summarize measurable engineering results. Avoid claiming scientific accuracy beyond the available evaluation fixtures and human review process.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Demo readiness and operations", "Local-first setup with a documented backup path")
    add_card(slide, "Primary demo", "Run Postgres/pgvector, ingest approved PDFs, start FastAPI + Next.js, and use Ollama or OpenAI for grounded Ask.", x=0.8, y=2.15, w=5.65, h=2.0)
    add_card(slide, "Backup demo", "Use repository screenshots and this deck if network, model, Docker, or showcase equipment is unavailable.", x=6.85, y=2.15, w=5.65, h=2.0, accent=GOLD)
    add_bullets(slide, ["Documented clean-machine setup and troubleshooting", "PDF fetch/ingest commands and provenance-preserving index", "Static corpus browsing remains available without the model provider"], x=1.0, y=4.75, w=11.2, h=1.4, size=17)
    add_footer(slide, 10)
    add_notes(slide, "Explain the live-demo sequence and backup plan. These committed artifacts support evaluation and handoff if the live environment is unavailable.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Known limitations and future work", "The MVP is intentionally bounded")
    add_bullets(slide, ["Local deployment only; no public production URL", "Grounded Ask requires an indexed corpus plus configured OpenAI or Ollama provider", "Scientific outputs still require human review", "Corpus focuses on microgravity and skeletal muscle", "Future: stronger retrieval benchmarks, broader curated corpus, reviewer workflow, and production deployment", "Graph extraction remains experimental; no graph database in the MVP"], x=1.0, y=2.15, w=11.0, h=4.3, size=20)
    add_footer(slide, 11)
    add_notes(slide, "Be candid about limitations. Emphasize that local-first deployment and a controlled topic were conscious scope decisions, not hidden gaps.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, dark=True)
    add_logo(slide, x=0.85, y=0.7, width=4.0)
    add_title(slide, "Trustworthy evidence, one passage at a time", "A practical foundation for citation-first space-biology research", dark=True)
    add_bullets(slide, ["Controlled corpus", "Grounded retrieval", "Validated citations", "Transparent limitations"], x=0.95, y=2.65, w=5.0, h=2.5, dark=True, size=24)
    contact = slide.shapes.add_textbox(Inches(7.2), Inches(4.7), Inches(5.2), Inches(1.2))
    p = contact.text_frame.paragraphs[0]
    p.text = "John Hernandez\njherna65@fau.edu\nGitHub: bluefate/spacebio-evidence-engine"
    p.font.name = "Aptos"
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.RIGHT
    add_footer(slide, 12, dark=True)
    add_notes(slide, "Close by returning to the core value: scientific synthesis that stays linked to inspectable evidence. Invite repository-based review and questions.")

    prs.save(FINAL / "Space_Biology_Evidence_Engine_Pitch_Deck.pptx")


def build_summary() -> None:
    prs = new_deck()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, dark=True)
    add_logo(slide, x=0.55, y=0.35, width=3.2)

    title = slide.shapes.add_textbox(Inches(0.65), Inches(1.22), Inches(7.2), Inches(0.65))
    p = title.text_frame.paragraphs[0]
    p.text = "Citation-first evidence for\nspace-biology research"
    p.font.name = "Aptos Display"; p.font.size = Pt(25); p.font.bold = True; p.font.color.rgb = WHITE
    sub = slide.shapes.add_textbox(Inches(0.68), Inches(2.05), Inches(7.7), Inches(0.52))
    p = sub.text_frame.paragraphs[0]
    p.text = "Local-first RAG for 23 approved microgravity and skeletal-muscle publications."
    p.font.name = "Aptos"; p.font.size = Pt(13); p.font.color.rgb = RGBColor(205, 230, 240)

    add_card(slide, "Who it serves", "Researchers, students, educators, and corpus maintainers who need inspectable evidence—not unsupported model recall.", x=0.65, y=2.75, w=3.85, h=1.45)
    add_card(slide, "Core capabilities", "Semantic + keyword search; grounded Q&A; passage citations; provenance; evidence sufficiency; study comparison.", x=4.72, y=2.75, w=3.85, h=1.45, accent=GOLD)
    add_card(slide, "Current status", "MVP complete and demo-ready locally. Code, CI, docs, ingestion, retrieval, Ask, citation UI, and evaluations are complete.", x=0.65, y=4.45, w=3.85, h=1.55, accent=BLUE)
    add_card(slide, "Demo readiness", "Docker + indexed PDFs + Ollama/OpenAI support the live app; committed screenshots and presentation materials provide a reliable backup.", x=4.72, y=4.45, w=3.85, h=1.55, accent=CYAN)
    add_screenshot(slide, "01-home.png", x=8.85, y=1.55, w=4.05, h=2.82)
    add_screenshot(slide, "02-ask.png", x=8.85, y=4.55, w=4.05, h=2.25)

    contact = slide.shapes.add_textbox(Inches(0.72), Inches(6.35), Inches(7.7), Inches(0.6))
    p = contact.text_frame.paragraphs[0]
    p.text = "John Hernandez  |  jherna65@fau.edu  |  Sponsor: FAU AI HootCamp  |  Repository: github.com/bluefate/spacebio-evidence-engine"
    p.font.name = "Aptos"; p.font.size = Pt(10); p.font.color.rgb = WHITE
    add_notes(slide, "One-page project summary covering the intended users, core capabilities, final implementation status, and demo readiness.")
    prs.save(FINAL / "Space_Biology_Evidence_Engine_One_Page_Summary.pptx")


if __name__ == "__main__":
    FINAL.mkdir(parents=True, exist_ok=True)
    build_pitch()
    build_summary()
